from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import io
from docx.shared import Inches
import httpx
import re
import os


class DocumentService:

    @staticmethod
    def _get_image_stream(url: str):
        try:
            with httpx.Client() as client:
                resp = client.get(url)
                if resp.status_code == 200:
                    return io.BytesIO(resp.content)
        except Exception as e:
            print(f"Failed to download image: {e}")
        return None

    @staticmethod
    def _clean_text(text: str):
        """Removes markdown symbols like ** and ## for cleaner output"""
        if not text: return ""
        clean = re.sub(r'\*\*|##', '', text)
        return clean.strip()

    @staticmethod
    def generate_word(title: str, sections: list) -> bytes:
        doc = Document()
        #Title
        doc.add_heading(title, 0)
        for section in sections:
            #Clean title
            clean_title = DocumentService._clean_text(section.get("title", ""))
            doc.add_heading(clean_title, level=1)
            
            if section.get("content"):
                #Clean content
                clean_content = DocumentService._clean_text(section.get("content", ""))
                doc.add_paragraph(clean_content)
            
            image_url = section.get("image_url")
            if image_url:
                image_stream = DocumentService._get_image_stream(image_url)
                if image_stream:
                    try:
                        doc.add_picture(image_stream, width=Inches(5))
                    except Exception as e:
                        print(f"Error adding picture to Word: {e}")
                        doc.add_paragraph("[Image could not be loaded]")

            doc.add_paragraph()
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    @staticmethod
    def generate_powerpoint(title: str, sections: list, template: str = "default") -> bytes:
        base_path = os.path.dirname(os.path.dirname(__file__))
        template_path = os.path.join(base_path, "templates", f"{template}.pptx")

        #Load Presentation
        if os.path.exists(template_path):
            prs = Presentation(template_path)
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
        else:
            print(f"Template not found at {template_path}, using blank.")
            prs = Presentation() 
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        
        #Title Slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        #Content Slides
        for section in sections:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            

            if slide.shapes.title:
                clean_title = DocumentService._clean_text(section.get("title", ""))
                slide.shapes.title.text = clean_title


            if section.get("content"):

                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    if hasattr(body, "text_frame"):
                        body.text_frame.clear() 
                        
                        content_raw = section["content"].split("\n")
                        for line in content_raw:
                            line = line.strip()
                            if not line: continue
                            
                            clean_line = re.sub(r'^[-*•]\s*', '', DocumentService._clean_text(line))
                            
                            p = body.text_frame.add_paragraph()
                            p.text = clean_line
                            p.level = 0

                            p.font.size = Pt(12)  
                            p.font.name = "Arial" 
                            p.space_after = Pt(6) 
                            


            image_url = section.get("image_url")
            if image_url:
                image_stream = DocumentService._get_image_stream(image_url)
                if image_stream:
                    try:

                        slide.shapes.add_picture(
                            image_stream, left=Inches(8.5), top=Inches(2), width=Inches(4)
                        )
                    except Exception:
                        pass

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

document_service = DocumentService()